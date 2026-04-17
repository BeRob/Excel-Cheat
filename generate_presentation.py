import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- Konfiguration ---
OUTPUT_FILE = "Messanwendung_Praesentation.pptx"
SCREENSHOT_DIR = "screenshots"

# QuestAlpha Farben (aus theme.py abgeleitet)
COLOR_PRIMARY = RGBColor(0x1C, 0x1C, 0x1C)  # Anthrazit
COLOR_ACCENT = RGBColor(0x00, 0x66, 0xCC)   # Blau
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GREY = RGBColor(0xF8, 0xF9, 0xFA)

def create_presentation():
    prs = Presentation()

    # --- Hilfsfunktionen ---
    def set_background(slide, color=COLOR_WHITE):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(prs, title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.name = 'Segoe UI'
        
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
        subtitle.text_frame.paragraphs[0].font.name = 'Segoe UI'

    def add_bullet_slide(prs, title_text, items):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
        title.text_frame.paragraphs[0].font.name = 'Segoe UI'
        title.text_frame.paragraphs[0].font.bold = True
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() 

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.name = 'Segoe UI'
            p.space_after = Pt(14)

    def add_image_slide(prs, title_text, image_paths, captions):
        slide_layout = prs.slide_layouts[5] # Title Only (Custom layout essentially)
        slide = prs.slides.add_slide(slide_layout)
        set_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
        title.text_frame.paragraphs[0].font.name = 'Segoe UI'
        title.text_frame.paragraphs[0].font.bold = True

        # Layout calculation for 2 images side-by-side
        left_margin = Inches(0.5)
        top_margin = Inches(2.0)
        img_width = Inches(4.2)
        gap = Inches(0.5)

        for i, img_path in enumerate(image_paths):
            full_path = os.path.join(SCREENSHOT_DIR, img_path)
            if os.path.exists(full_path):
                left = left_margin + (i * (img_width + gap))
                pic = slide.shapes.add_picture(full_path, left, top_margin, width=img_width)
                
                # Add Caption
                textbox = slide.shapes.add_textbox(left, top_margin + Inches(3.2), img_width, Inches(0.5))
                p = textbox.text_frame.add_paragraph()
                p.text = captions[i]
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(12)
                p.font.color.rgb = COLOR_PRIMARY

    # --- Folien Inhalte ---

    # 1. Titel
    add_title_slide(prs, "Messanwendung (Easy Mode)", "Effiziente Datenerfassung & Qualitätssicherung\nQuestAlpha")

    # 2. Warum diese Anwendung? (Herausforderung vs Lösung)
    add_bullet_slide(prs, "Ausgangslage & Zielsetzung", [
        "Herausforderung: Manuelle Bearbeitung von Excel-Listen ist fehleranfällig.",
        "Risiko: Versehentliches Löschen von Formeln oder Überschreiben alter Werte.",
        "Problem: Fehlende Nachvollziehbarkeit (Wer hat wann was geändert?).",
        "Lösung: Eine geführte Desktop-Anwendung als sichere Schnittstelle.",
        "Ziel: Datenintegrität gewährleisten und Prozess vereinfachen."
    ])

    # 3. Der Nutzen
    add_bullet_slide(prs, "Vorteile der Anwendung", [
        "✅ Sicherheit: Schreibschutz für bestehende Daten (Append-Only).",
        "✅ Qualität: Automatische Validierung der Eingaben (Zahlenformate).",
        "✅ Audit-Trail: Lückenlose Protokollierung aller Aktionen.",
        "✅ Effizienz: Kontext (Charge, Auftrag) muss nur einmal gesetzt werden.",
        "✅ Benutzerfreundlichkeit: Klare, moderne Oberfläche im QuestAlpha-Design."
    ])

    # 4. Funktionsüberblick
    add_bullet_slide(prs, "Funktionsweise & Workflow", [
        "1. Login: Sicherer Zugang via Passwort oder QR-Code Scan.",
        "2. Dateiwahl: Auswahl der prozessspezifischen Excel-Datei.",
        "3. Konfiguration: Festlegen von Kontextdaten (z.B. Chargennummer).",
        "4. Erfassung: Eingabe der Messwerte in ein generiertes Formular.",
        "5. Speichern: Daten werden automatisch validiert und angehängt."
    ])

    # 5. Einblick I (Start)
    add_image_slide(prs, "Einfacher Einstieg", 
                   ["1_login.png", "2_datei.png"],
                   ["Login Screen", "Dateiauswahl"])

    # 6. Einblick II (Erfassung)
    add_image_slide(prs, "Datenerfassung", 
                   ["4_kontext.png", "5_formular.png"],
                   ["Kontext setzen", "Messwerte eingeben"])

    prs.save(OUTPUT_FILE)
    print(f"Präsentation erfolgreich erstellt: {OUTPUT_FILE}")

if __name__ == "__main__":
    create_presentation()
